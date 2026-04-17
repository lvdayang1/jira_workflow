# -*- coding: utf-8 -*-
"""
附件读取器 - 支持多种文件格式的文本提取
"""
import os
import logging
from pathlib import Path

logger = logging.getLogger(__name__)


class AttachmentReader:
    """附件读取器基类"""

    def can_read(self, filepath: str) -> bool:
        """检查是否支持读取此文件"""
        raise NotImplementedError

    def read(self, filepath: str) -> str:
        """读取文件内容并返回文本"""
        raise NotImplementedError


class PPTXReader(AttachmentReader):
    """PPTX 文件读取器"""

    def can_read(self, filepath: str) -> bool:
        return filepath.lower().endswith('.pptx')

    def read(self, filepath: str) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            text_parts = []

            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                if slide_texts:
                    text_parts.append(f"=== 第{i+1}页 ===")
                    text_parts.append('\n'.join(slide_texts))

            return '\n\n'.join(text_parts) if text_parts else ""
        except Exception as e:
            logger.error(f"读取PPTX失败 {filepath}: {e}")
            return f"[PPTX读取失败: {e}]"


class DOCXReader(AttachmentReader):
    """DOCX 文件读取器"""

    def can_read(self, filepath: str) -> bool:
        return filepath.lower().endswith('.docx')

    def read(self, filepath: str) -> str:
        try:
            from docx import Document
            doc = Document(filepath)
            paragraphs = []

            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text.strip())

            # 读取表格
            for table in doc.tables:
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        paragraphs.append(' | '.join(cells))

            return '\n'.join(paragraphs) if paragraphs else ""
        except Exception as e:
            logger.error(f"读取DOCX失败 {filepath}: {e}")
            return f"[DOCX读取失败: {e}]"


class ExcelReader(AttachmentReader):
    """Excel 文件读取器 (xlsx/xls)"""

    def can_read(self, filepath: str) -> bool:
        return filepath.lower().endswith(('.xlsx', '.xls'))

    def read(self, filepath: str) -> str:
        try:
            import openpyxl
            wb = openpyxl.load_workbook(filepath, data_only=True)
            text_parts = []

            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                text_parts.append(f"=== Sheet: {sheet_name} ===")

                rows_data = []
                for row in sheet.iter_rows(max_row=100, values_only=True):
                    row_values = [str(cell) if cell is not None else '' for cell in row]
                    if any(v.strip() for v in row_values):
                        rows_data.append(' | '.join(row_values))

                text_parts.append('\n'.join(rows_data[:50]))  # 限制行数

            return '\n\n'.join(text_parts) if text_parts else ""
        except Exception as e:
            logger.error(f"读取Excel失败 {filepath}: {e}")
            return f"[Excel读取失败: {e}]"


class ImageReader(AttachmentReader):
    """图片文件读取器 - 提取基本信息和OCR（如果可用）"""

    def can_read(self, filepath: str) -> bool:
        return filepath.lower().endswith(('.png', '.jpg', '.jpeg', '.gif', '.bmp', '.webp'))

    def read(self, filepath: str) -> str:
        try:
            from PIL import Image
            img = Image.open(filepath)
            info = [
                f"文件名: {os.path.basename(filepath)}",
                f"图片尺寸: {img.width}x{img.height}",
                f"图片格式: {img.format}",
            ]

            # 尝试OCR（如果安装了pytesseract）
            try:
                import pytesseract
                text = pytesseract.image_to_string(img, lang='chi_sim+eng')
                if text.strip():
                    info.append("OCR识别内容:")
                    info.append(text.strip())
            except ImportError:
                info.append("[提示: 未安装OCR库，无法提取图片文字]")

            return '\n'.join(info)
        except Exception as e:
            logger.error(f"读取图片失败 {filepath}: {e}")
            return f"[图片读取失败: {e}]"


class TextFileReader(AttachmentReader):
    """文本文件读取器"""

    def can_read(self, filepath: str) -> bool:
        return filepath.lower().endswith(('.txt', '.md', '.log', '.csv', '.json', '.xml', '.html'))

    def read(self, filepath: str) -> str:
        try:
            encodings = ['utf-8', 'gbk', 'gb2312', 'latin1']
            for encoding in encodings:
                try:
                    with open(filepath, 'r', encoding=encoding) as f:
                        content = f.read()
                    return content[:5000]  # 限制长度
                except UnicodeDecodeError:
                    continue
            return "[文本文件编码不支持]"
        except Exception as e:
            logger.error(f"读取文本文件失败 {filepath}: {e}")
            return f"[文本文件读取失败: {e}]"


class PDFReader(AttachmentReader):
    """PDF 文件读取器"""

    def can_read(self, filepath: str) -> bool:
        return filepath.lower().endswith('.pdf')

    def read(self, filepath: str) -> str:
        try:
            # 尝试使用PyPDF2
            import PyPDF2
            with open(filepath, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                text_parts = []
                for i, page in enumerate(reader.pages[:10]):  # 限制页数
                    text = page.extract_text()
                    if text:
                        text_parts.append(f"=== 第{i+1}页 ===")
                        text_parts.append(text)
                return '\n\n'.join(text_parts) if text_parts else "[PDF无文本内容]"
        except ImportError:
            try:
                # 尝试使用pdfplumber
                import pdfplumber
                with pdfplumber.open(filepath) as pdf:
                    text_parts = []
                    for i, page in enumerate(pdf.pages[:10]):
                        text = page.extract_text()
                        if text:
                            text_parts.append(f"=== 第{i+1}页 ===")
                            text_parts.append(text)
                    return '\n\n'.join(text_parts) if text_parts else "[PDF无文本内容]"
            except ImportError:
                return "[PDF读取需要安装 PyPDF2 或 pdfplumber]"
        except Exception as e:
            logger.error(f"读取PDF失败 {filepath}: {e}")
            return f"[PDF读取失败: {e}]"


class UnsupportedReader(AttachmentReader):
    """不支持的文件格式"""

    def can_read(self, filepath: str) -> bool:
        return True  # 作为最后兜底

    def read(self, filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()
        return f"[不支持的文件格式: {ext}]"


def get_all_readers():
    """获取所有文件读取器"""
    return [
        PPTXReader(),
        DOCXReader(),
        ExcelReader(),
        PDFReader(),
        ImageReader(),
        TextFileReader(),
        UnsupportedReader(),
    ]


def read_attachment(filepath: str) -> str:
    """
    读取附件文件内容

    Args:
        filepath: 文件路径

    Returns:
        文件文本内容
    """
    if not os.path.exists(filepath):
        return f"[文件不存在: {filepath}]"

    for reader in get_all_readers():
        if reader.can_read(filepath):
            return reader.read(filepath)

    return "[无法读取文件]"
